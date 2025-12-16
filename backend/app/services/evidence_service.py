# backend/services/evidence_processor.py
import asyncio
import concurrent.futures
import hashlib
import logging
import mimetypes
import os
import tempfile
import time
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

try:
    from app.services.ai_service import ai_service
    from app.services.search_service import evidence_search_index

    vector_store = ai_service.vector_store
except ImportError:
    from app.services.ai_service import ai_service
    from app.services.search_service import evidence_search_index

    vector_store = ai_service.vector_store

logger = logging.getLogger(__name__)


@dataclass
class ProcessingResult:
    file_id: str
    file_path: str
    file_type: str
    size_bytes: int
    processing_time: float
    extracted_text: str = ""
    key_entities: List[Dict[str, Any]] = None
    sentiment_score: float = 0.0
    quality_score: float = 0.0
    error: str = None
    metadata: Dict[str, Any] = None

    def __post_init__(self):
        if self.key_entities is None:
            self.key_entities = []
        if self.metadata is None:
            self.metadata = {}


class EvidenceProcessor:
    """
    Optimized evidence processing pipeline with parallel processing and memory management
    """

    def __init__(self, max_workers: int = 4, temp_dir: str = None):
        self.max_workers = max_workers
        self.temp_dir = temp_dir or tempfile.gettempdir()
        self.executor = concurrent.futures.ThreadPoolExecutor(max_workers=max_workers)

        # Processing capabilities
        self.supported_types = {
            "image": [
                "image/jpeg",
                "image/png",
                "image/tiff",
                "image/bmp",
                "image/gif",
                "image/webp",
            ],
            "document": [
                "application/pdf",
                "application/msword",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/vnd.ms-excel",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-powerpoint",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ],
            "text": [
                "text/plain",
                "text/csv",
                "text/html",
                "text/markdown",
                "application/json",
            ],
            "audio": [
                "audio/mpeg",
                "audio/wav",
                "audio/ogg",
                "audio/mp4",
                "audio/webm",
            ],
            "video": [
                "video/mp4",
                "video/avi",
                "video/mov",
                "video/wmv",
                "video/webm",
                "video/ogg",
            ],
            "archive": [
                "application/zip",
                "application/x-rar-compressed",
                "application/x-7z-compressed",
            ],
        }

        # Performance metrics
        self.metrics = {
            "total_processed": 0,
            "total_processing_time": 0.0,
            "errors": 0,
            "by_type": {},
        }

    async def process_files_batch(
        self, file_paths: List[str], options: Dict[str, Any] = None
    ) -> List[ProcessingResult]:
        """
        Process multiple files in parallel with optimized resource usage
        """
        if not file_paths:
            return []

        options = options or {}
        batch_start = time.time()

        # Create processing tasks
        tasks = []
        for file_path in file_paths:
            task = self._process_single_file_async(file_path, options)
            tasks.append(task)

        # Process in batches to avoid overwhelming the system
        batch_size = min(self.max_workers * 2, len(tasks))
        results = []

        for i in range(0, len(tasks), batch_size):
            batch_tasks = tasks[i : i + batch_size]
            batch_results = await asyncio.gather(*batch_tasks, return_exceptions=True)

            for result in batch_results:
                if isinstance(result, Exception):
                    logger.error(f"Batch processing error: {result}")
                    results.append(
                        ProcessingResult(
                            file_id="",
                            file_path="",
                            file_type="",
                            size_bytes=0,
                            processing_time=0.0,
                            error=str(result),
                        )
                    )
                else:
                    results.append(result)

            # Small delay between batches to prevent resource exhaustion
            if i + batch_size < len(tasks):
                await asyncio.sleep(0.1)

        # Update metrics
        batch_time = time.time() - batch_start
        self._update_batch_metrics(results, batch_time)

        logger.info(f"Processed {len(file_paths)} files in {batch_time:.2f}s")
        return results

    async def _process_single_file_async(
        self, file_path: str, options: Dict[str, Any]
    ) -> ProcessingResult:
        """
        Process a single file asynchronously
        """
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            self.executor, self._process_single_file_sync, file_path, options
        )

    def _process_single_file_sync(
        self, file_path: str, options: Dict[str, Any]
    ) -> ProcessingResult:
        """
        Process a single file synchronously (runs in thread pool)
        """
        start_time = time.time()
        file_id = self._generate_file_id(file_path)

        try:
            # Validate file
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")

            file_size = os.path.getsize(file_path)
            mime_type = self._detect_mime_type(file_path)

            # Check file size limits
            max_size = options.get("max_file_size", 50 * 1024 * 1024)  # 50MB default
            if file_size > max_size:
                raise ValueError(f"File too large: {file_size} bytes (max: {max_size})")

            # Process based on file type
            result = ProcessingResult(
                file_id=file_id,
                file_path=file_path,
                file_type=mime_type,
                size_bytes=file_size,
                processing_time=0.0,
            )

            if mime_type in self.supported_types["image"]:
                self._process_image(file_path, result, options)
            elif mime_type in self.supported_types["document"]:
                self._process_document(file_path, result, options)
            elif mime_type in self.supported_types["text"]:
                self._process_text(file_path, result, options)
            elif mime_type in self.supported_types["audio"]:
                self._process_audio(file_path, result, options)
            elif mime_type in self.supported_types["video"]:
                self._process_video(file_path, result, options)
            elif mime_type in [
                "application/vnd.ms-excel",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "text/csv",
            ]:
                self._process_spreadsheet(file_path, result, options)
            elif mime_type in [
                "application/vnd.ms-powerpoint",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ]:
                self._process_presentation(file_path, result, options)
            elif mime_type in self.supported_types["archive"]:
                self._process_archive(file_path, result, options)
            else:
                result.error = f"Unsupported file type: {mime_type}"

            result.processing_time = time.time() - start_time

            # Index the processed evidence for search
            if not result.error:
                try:
                    processing_dict = {
                        "extracted_text": result.extracted_text,
                        "key_entities": result.key_entities,
                        "metadata": result.metadata,
                        "file_type": result.file_type,
                        "quality_score": result.quality_score,
                        "sentiment_score": result.sentiment_score,
                    }
                    evidence_search_index.index_evidence(
                        result.file_id, result.file_path, processing_dict
                    )

                    # Add to vector store for semantic search
                    content_to_embed = result.extracted_text or result.file_path
                    if content_to_embed.strip():
                        vector_store.add_document(
                            document_id=result.file_id,
                            content=content_to_embed,
                            metadata={
                                "file_path": result.file_path,
                                "file_type": result.file_type,
                                "quality_score": result.quality_score,
                                "sentiment_score": result.sentiment_score,
                            },
                        )

                except Exception as index_error:
                    logger.warning(
                        f"Failed to index evidence {result.file_id}: {index_error}"
                    )

            return result

        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            return ProcessingResult(
                file_id=file_id,
                file_path=file_path,
                file_type="",
                size_bytes=0,
                processing_time=time.time() - start_time,
                error=str(e),
            )

    def _process_image(
        self, file_path: str, result: ProcessingResult, options: Dict[str, Any]
    ):
        """Process image files with OCR and analysis"""
        try:
            # Import here to avoid import errors if PIL is not available
            import cv2
            import numpy as np
            import pytesseract
            from PIL import Image

            # Open image
            image = Image.open(file_path)

            # Basic image analysis
            result.metadata.update(
                {
                    "width": image.width,
                    "height": image.height,
                    "format": image.format,
                    "mode": image.mode,
                }
            )

            # OCR processing
            if options.get("enable_ocr", True):
                try:
                    # Convert PIL to numpy array for OpenCV processing
                    if image.mode != "RGB":
                        image = image.convert("RGB")

                    # Enhance image for better OCR
                    opencv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)

                    # Apply preprocessing for better OCR accuracy
                    gray = cv2.cvtColor(opencv_image, cv2.COLOR_BGR2GRAY)
                    _, threshold = cv2.threshold(
                        gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU
                    )

                    # OCR
                    text = pytesseract.image_to_string(threshold, lang="eng")
                    result.extracted_text = text.strip()

                    # Extract key entities (simplified - could use NLP models)
                    result.key_entities = self._extract_entities_from_text(text)

                except Exception as ocr_error:
                    logger.warning(
                        f"OCR processing failed for {file_path}: {ocr_error}"
                    )
                    result.metadata["ocr_error"] = str(ocr_error)

            # Image forensics analysis
            forensics_result = self._analyze_image_forensics(opencv_image, image)
            result.metadata.update(forensics_result)

            # Quality assessment
            result.quality_score = self._assess_image_quality(image)

            # Sentiment analysis (if text was extracted)
            if result.extracted_text:
                result.sentiment_score = self._analyze_sentiment(result.extracted_text)

        except ImportError as e:
            result.error = f"Image processing dependencies not available: {e}"
        except Exception as e:
            result.error = f"Image processing failed: {e}"

    def _analyze_image_forensics(
        self, opencv_image: "np.ndarray", pil_image: "Image"
    ) -> Dict[str, Any]:
        """Analyze image for forensic indicators with scoring"""
        forensics = {}
        manipulation_score = 0.0
        authenticity_score = 100.0
        forensic_indicators = []

        try:
            import cv2
            import numpy as np
            from PIL import Image

            # Error Level Analysis (ELA) - detects image manipulation
            ela_result = self._error_level_analysis(pil_image)
            forensics.update(ela_result)
            if ela_result.get("ela_score", 0) > 15:
                forensic_indicators.append(
                    "High error level analysis score - possible manipulation"
                )
                manipulation_score += 25
                authenticity_score -= 20

            # Noise analysis
            noise_result = self._analyze_image_noise(opencv_image)
            forensics.update(noise_result)
            if noise_result.get("noise_level", 0) > 0.8:  # Arbitrary threshold example
                forensic_indicators.append("High noise level detected")
                # Adjust scores as appropriate

            # Metadata analysis
            metadata_result = self._analyze_image_metadata(pil_image)
            forensics.update(metadata_result)
            if metadata_result.get("suspicious_software"):
                forensic_indicators.append(
                    f"Suspicious software detected: {metadata_result['suspicious_software']}"
                )
                manipulation_score += 15
                authenticity_score -= 10
            if not metadata_result.get("date_consistency", True):
                 forensic_indicators.append("Inconsistent dates in metadata")
                 manipulation_score += 10
                 authenticity_score -= 5

            # Compression artifacts
            compression_result = self._detect_compression_artifacts(opencv_image)
            forensics.update(compression_result)
            if compression_result.get("likely_compressed"):
                 forensic_indicators.append("Likely re-compressed")

            # Clone detection (basic)
            clone_result = self._detect_clone_regions(opencv_image)
            forensics.update(clone_result)
            if clone_result.get("clone_regions_detected"):
                forensic_indicators.append("Clone regions detected")
                manipulation_score += 30
                authenticity_score -= 25

            # Cap scores
            manipulation_score = min(manipulation_score, 100.0)
            authenticity_score = max(authenticity_score, 0.0)

            forensics.update({
                "manipulation_score": manipulation_score,
                "authenticity_score": authenticity_score,
                "forensic_indicators": forensic_indicators
            })

        except Exception as e:
            forensics["forensics_error"] = str(e)

        return forensics

    def _error_level_analysis(self, image: "Image") -> Dict[str, Any]:
        """Perform Error Level Analysis to detect manipulation"""
        try:
            import cv2
            import numpy as np

            # Save image with high quality JPEG
            temp_path = os.path.join(
                self.temp_dir, f"temp_ela_{hash(image.tobytes())}.jpg"
            )
            image.save(temp_path, "JPEG", quality=95)

            # Reload and compare
            reloaded = Image.open(temp_path)
            original_array = np.array(image)
            reloaded_array = np.array(reloaded)

            # Calculate difference
            if original_array.shape == reloaded_array.shape:
                diff = np.abs(
                    original_array.astype(np.int16) - reloaded_array.astype(np.int16)
                )
                ela_score = np.mean(diff)

                # Clean up
                os.unlink(temp_path)

                return {
                    "ela_score": float(ela_score),
                    "manipulation_likelihood": "high" if ela_score > 10 else "low",
                }
            else:
                os.unlink(temp_path)
                return {"ela_score": 0.0, "manipulation_likelihood": "unknown"}

        except Exception as e:
            return {"ela_error": str(e)}

    def _analyze_image_noise(self, image: "np.ndarray") -> Dict[str, Any]:
        """Analyze image noise patterns"""
        try:
            import cv2

            # Convert to grayscale if needed
            if len(image.shape) == 3:
                gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            else:
                gray = image

            # Calculate noise using Laplacian variance
            laplacian_var = cv2.Laplacian(gray, cv2.CV_64F).var()

            # Estimate noise level
            noise_level = 1.0 / (1.0 + laplacian_var / 1000.0)

            return {
                "noise_level": float(noise_level),
                "sharpness_score": float(laplacian_var),
            }

        except Exception as e:
            return {"noise_analysis_error": str(e)}

    def _analyze_image_metadata(self, image: "Image") -> Dict[str, Any]:
        """Analyze image metadata for forensic clues"""
        metadata = {}

        try:
            # Extract EXIF data
            exif_data = image._getexif()
            if exif_data:
                metadata["has_exif"] = True

                # Check for suspicious metadata
                suspicious_software = ["photoshop", "gimp", "paint"]
                software = exif_data.get(305, "").lower()  # Software tag
                if any(s in software for s in suspicious_software):
                    metadata["suspicious_software"] = software

                # Check creation date vs modification date
                date_original = exif_data.get(36867)  # DateTimeOriginal
                date_digitized = exif_data.get(36868)  # DateTimeDigitized

                if date_original and date_digitized:
                    metadata["date_consistency"] = date_original == date_digitized

            else:
                metadata["has_exif"] = False

        except Exception as e:
            metadata["metadata_error"] = str(e)

        return metadata

    def _detect_compression_artifacts(self, image: "np.ndarray") -> Dict[str, Any]:
        """Detect JPEG compression artifacts"""
        try:
            import cv2

            # Convert to grayscale
            if len(image.shape) == 3:
                gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            else:
                gray = image

            # High-pass filter to detect blocking artifacts
            kernel = np.array([[-1, -1, -1], [-1, 8, -1], [-1, -1, -1]])
            high_pass = cv2.filter2D(gray, -1, kernel)

            # Calculate artifact score
            artifact_score = np.mean(np.abs(high_pass))

            return {
                "compression_artifacts": float(artifact_score),
                "likely_compressed": artifact_score > 50,
            }

        except Exception as e:
            return {"compression_analysis_error": str(e)}

    def _detect_clone_regions(self, image: "np.ndarray") -> Dict[str, Any]:
        """Basic clone region detection"""
        try:
            import cv2

            # Convert to grayscale
            if len(image.shape) == 3:
                gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            else:
                gray = image

            # Use template matching to find similar regions
            # This is a simplified approach - real forensic tools use more sophisticated methods
            height, width = gray.shape

            # Check for identical blocks (very basic clone detection)
            clone_detected = False
            block_size = 32

            for y in range(0, height - block_size, block_size):
                for x in range(0, width - block_size, block_size):
                    block = gray[y : y + block_size, x : x + block_size]

                    # Look for identical blocks elsewhere
                    for y2 in range(y + block_size, height - block_size, block_size):
                        for x2 in range(0, width - block_size, block_size):
                            if x == x2 and y == y2:
                                continue

                            block2 = gray[y2 : y2 + block_size, x2 : x2 + block_size]
                            if np.array_equal(block, block2):
                                clone_detected = True
                                break
                        if clone_detected:
                            break
                    if clone_detected:
                        break
                if clone_detected:
                    break

            return {
                "clone_regions_detected": clone_detected,
                "clone_detection_method": "basic_block_comparison",
            }

        except Exception as e:
            return {"clone_detection_error": str(e)}

    def _process_document(
        self, file_path: str, result: ProcessingResult, options: Dict[str, Any]
    ):
        """Process document files (PDF, DOCX, etc.)"""
        try:
            mime_type = result.file_type

            if mime_type == "application/pdf":
                self._process_pdf(file_path, result, options)
            elif mime_type in [
                "application/msword",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ]:
                self._process_docx(file_path, result, options)
            else:
                result.error = f"Unsupported document type: {mime_type}"

        except Exception as e:
            result.error = f"Document processing failed: {e}"

    def _process_pdf(
        self, file_path: str, result: ProcessingResult, options: Dict[str, Any]
    ):
        """Process PDF files"""
        try:
            # Import here to avoid import errors
            import fitz  # PyMuPDF

            doc = fitz.open(file_path)

            # Extract text from all pages
            text = ""
            for page in doc:
                text += page.get_text() + "\n"

            result.extracted_text = text.strip()
            result.metadata.update(
                {
                    "pages": len(doc),
                    "title": doc.metadata.get("title", ""),
                    "author": doc.metadata.get("author", ""),
                    "subject": doc.metadata.get("subject", ""),
                }
            )

            # Extract entities
            result.key_entities = self._extract_entities_from_text(text)

            # Quality assessment based on text length and structure
            result.quality_score = min(
                1.0, len(text) / 10000
            )  # Normalize by expected content

            # Sentiment analysis
            result.sentiment_score = self._analyze_sentiment(text)

            doc.close()

        except ImportError:
            result.error = "PDF processing requires PyMuPDF (fitz)"
        except Exception as e:
            result.error = f"PDF processing failed: {e}"

    def _process_docx(
        self, file_path: str, result: ProcessingResult, options: Dict[str, Any]
    ):
        """Process DOCX files"""
        try:
            # Import here to avoid import errors
            from docx import Document

            doc = Document(file_path)

            # Extract text from all paragraphs
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"

            result.extracted_text = text.strip()
            result.metadata.update(
                {"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}
            )

            # Extract entities
            result.key_entities = self._extract_entities_from_text(text)

            # Quality assessment
            result.quality_score = min(1.0, len(text) / 5000)

            # Sentiment analysis
            result.sentiment_score = self._analyze_sentiment(text)

        except ImportError:
            result.error = "DOCX processing requires python-docx"
        except Exception as e:
            result.error = f"DOCX processing failed: {e}"

    def _process_text(
        self, file_path: str, result: ProcessingResult, options: Dict[str, Any]
    ):
        """Process text files"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

            result.extracted_text = text
            result.metadata.update(
                {
                    "lines": len(text.split("\n")),
                    "characters": len(text),
                    "encoding": "utf-8",
                }
            )

            # Extract entities
            result.key_entities = self._extract_entities_from_text(text)

            # Quality assessment
            result.quality_score = min(1.0, len(text) / 1000)

            # Sentiment analysis
            result.sentiment_score = self._analyze_sentiment(text)

        except Exception as e:
            result.error = f"Text processing failed: {e}"

    def _process_audio(
        self, file_path: str, result: ProcessingResult, options: Dict[str, Any]
    ):
        """Process audio files for speech-to-text and analysis"""
        try:
            # Basic audio metadata extraction
            result.metadata.update(
                {
                    "media_type": "audio",
                    "processing_capabilities": ["metadata", "duration", "format"],
                }
            )

            # Placeholder for audio processing - would integrate with speech recognition
            # For now, just extract basic metadata
            try:
                import wave

                if file_path.lower().endswith(".wav"):
                    with wave.open(file_path, "rb") as wav_file:
                        duration = wav_file.getnframes() / wav_file.getframerate()
                        result.metadata.update(
                            {
                                "duration_seconds": duration,
                                "sample_rate": wav_file.getframerate(),
                                "channels": wav_file.getnchannels(),
                                "sample_width": wav_file.getsampwidth(),
                            }
                        )
            except ImportError:
                result.metadata["note"] = "Install wave module for WAV file analysis"

            # Speech-to-text placeholder (would integrate with Google Speech API, etc.)
            result.extracted_text = (
                "[Audio content - speech-to-text processing not implemented]"
            )
            result.quality_score = 0.5  # Placeholder quality score

        except Exception as e:
            result.error = f"Audio processing failed: {e}"

    def _process_video(
        self, file_path: str, result: ProcessingResult, options: Dict[str, Any]
    ):
        """Process video files for frame analysis and metadata"""
        try:
            result.metadata.update(
                {
                    "media_type": "video",
                    "processing_capabilities": ["metadata", "frames", "audio_track"],
                }
            )

            # Video metadata extraction (placeholder - would use OpenCV, etc.)
            result.metadata.update(
                {
                    "estimated_duration": "Unknown",  # Would extract from video
                    "resolution": "Unknown",  # Would extract width/height
                    "frame_rate": "Unknown",  # Would extract FPS
                    "has_audio": False,  # Would check for audio track
                }
            )

            # Frame analysis placeholder
            result.extracted_text = "[Video content - frame analysis not implemented]"
            result.key_entities = [
                {
                    "type": "video_metadata",
                    "confidence": 1.0,
                    "text": f"Video file: {os.path.basename(file_path)}",
                }
            ]
            result.quality_score = 0.6  # Placeholder quality score

        except Exception as e:
            result.error = f"Video processing failed: {e}"

    def _process_spreadsheet(
        self, file_path: str, result: ProcessingResult, options: Dict[str, Any]
    ):
        """Process Excel/CSV files"""
        try:
            import pandas as pd

            # Read spreadsheet
            if file_path.lower().endswith(".csv"):
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)

            # Extract text content
            text_content = []
            for col in df.columns:
                text_content.extend([str(val) for val in df[col].dropna()])

            result.extracted_text = "\n".join(text_content)

            # Metadata
            result.metadata.update(
                {
                    "rows": len(df),
                    "columns": len(df.columns),
                    "column_names": list(df.columns),
                    "data_types": {col: str(df[col].dtype) for col in df.columns},
                }
            )

            # Extract entities (financial data patterns)
            result.key_entities = self._extract_financial_entities(
                result.extracted_text
            )

            # Quality assessment
            result.quality_score = min(1.0, len(df) / 1000)  # Based on data volume

        except ImportError:
            result.error = "Spreadsheet processing requires pandas"
        except Exception as e:
            result.error = f"Spreadsheet processing failed: {e}"

    def _process_presentation(
        self, file_path: str, result: ProcessingResult, options: Dict[str, Any]
    ):
        """Process PowerPoint files"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_content = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)

            result.extracted_text = "\n\n".join(text_content)

            result.metadata.update(
                {"slides": len(prs.slides), "media_type": "presentation"}
            )

            # Extract entities
            result.key_entities = self._extract_entities_from_text(
                result.extracted_text
            )
            result.quality_score = min(1.0, len(text_content) / 100)

        except ImportError:
            result.error = "Presentation processing requires python-pptx"
        except Exception as e:
            result.error = f"Presentation processing failed: {e}"

    def _process_archive(
        self, file_path: str, result: ProcessingResult, options: Dict[str, Any]
    ):
        """Process archive files (ZIP, RAR, etc.)"""
        try:
            import zipfile

            with zipfile.ZipFile(file_path, "r") as zip_ref:
                file_list = zip_ref.namelist()
                file_info = []

                for info in zip_ref.filelist:
                    file_info.append(
                        {
                            "filename": info.filename,
                            "size": info.file_size,
                            "compressed_size": info.compress_size,
                            "date_time": str(info.date_time),
                        }
                    )

                result.metadata.update(
                    {
                        "archive_type": "zip",
                        "total_files": len(file_list),
                        "files": file_info,
                        "compression_ratio": sum(
                            f["compressed_size"] for f in file_info
                        )
                        / max(1, sum(f["size"] for f in file_info)),
                    }
                )

                # Extract text from text files in archive
                text_files = [
                    f
                    for f in file_list
                    if any(f.lower().endswith(ext) for ext in [".txt", ".md", ".csv"])
                ]
                extracted_texts = []

                for text_file in text_files[:5]:  # Limit to first 5 text files
                    try:
                        with zip_ref.open(text_file) as file:
                            content = file.read().decode("utf-8", errors="ignore")
                            extracted_texts.append(
                                f"=== {text_file} ===\n{content[:1000]}..."
                            )
                    except:
                        continue

                result.extracted_text = "\n\n".join(extracted_texts)
                result.quality_score = min(
                    1.0, len(file_list) / 50
                )  # Based on archive size

        except ImportError:
            result.error = "Archive processing requires zipfile (built-in)"
        except Exception as e:
            result.error = f"Archive processing failed: {e}"

    def _extract_financial_entities(self, text: str) -> List[Dict[str, Any]]:
        """Extract financial entities from text"""
        entities = []
        import re

        # Amount patterns
        amount_patterns = [
            r"\$[\d,]+\.?\d*",  # $1,234.56
            r"[\d,]+\.?\d*\s*(?:USD|dollars?|bucks?)",  # 1234 USD
        ]

        for pattern in amount_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                entities.append({"type": "amount", "text": match, "confidence": 0.9})

        # Account numbers (simplified)
        account_patterns = [
            r"\b\d{10,16}\b",  # 10-16 digit numbers (potential account numbers)
        ]

        for pattern in account_patterns:
            matches = re.findall(pattern, text)
            for match in matches:
                entities.append(
                    {"type": "account_number", "text": match, "confidence": 0.7}
                )

        return entities

    def _detect_mime_type(self, file_path: str) -> str:
        """Detect MIME type of file"""
        mime_type, _ = mimetypes.guess_type(file_path)
        return mime_type or "application/octet-stream"

    def _generate_file_id(self, file_path: str) -> str:
        """Generate unique file ID based on path and content"""
        file_path_hash = hashlib.md5(file_path.encode()).hexdigest()[:8]
        timestamp = str(int(time.time()))[-6:]  # Last 6 digits of timestamp
        return f"ev_{file_path_hash}_{timestamp}"

    def _extract_entities_from_text(self, text: str) -> List[Dict[str, Any]]:
        """Extract key entities from text (simplified implementation)"""
        entities = []

        # Simple pattern-based entity extraction
        import re

        # Email addresses
        emails = re.findall(
            r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b", text
        )
        for email in emails:
            entities.append({"type": "email", "value": email, "confidence": 0.9})

        # Phone numbers (basic pattern)
        phones = re.findall(r"\b\d{3}[-.]?\d{3}[-.]?\d{4}\b", text)
        for phone in phones:
            entities.append({"type": "phone", "value": phone, "confidence": 0.8})

        # Currency amounts
        amounts = re.findall(r"\$[\d,]+(?:\.\d{2})?", text)
        for amount in amounts:
            entities.append({"type": "currency", "value": amount, "confidence": 0.95})

        # Dates (basic pattern)
        dates = re.findall(r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b", text)
        for date in dates:
            entities.append({"type": "date", "value": date, "confidence": 0.7})

        return entities

    def _assess_image_quality(self, image) -> float:
        """Assess image quality for OCR and analysis"""
        try:
            # Basic quality metrics
            width, height = image.size
            total_pixels = width * height

            # Size factor
            size_score = min(1.0, total_pixels / (1920 * 1080))  # HD baseline

            # Color depth factor
            color_score = 1.0 if image.mode == "RGB" else 0.7

            # Overall quality score
            return (size_score + color_score) / 2

        except Exception:
            return 0.5

    def _analyze_sentiment(self, text: str) -> float:
        """Simple sentiment analysis (positive/negative scale)"""
        if not text:
            return 0.0

        # Simple keyword-based sentiment analysis
        positive_words = [
            "good",
            "excellent",
            "approved",
            "verified",
            "confirmed",
            "success",
        ]
        negative_words = [
            "bad",
            "fraud",
            "suspicious",
            "denied",
            "rejected",
            "error",
            "failed",
        ]

        text_lower = text.lower()
        positive_count = sum(1 for word in positive_words if word in text_lower)
        negative_count = sum(1 for word in negative_words if word in text_lower)

        total_sentiment_words = positive_count + negative_count

        if total_sentiment_words == 0:
            return 0.0

        # Normalize to -1 to 1 scale
        sentiment = (positive_count - negative_count) / total_sentiment_words
        return max(-1.0, min(1.0, sentiment))

    def _update_batch_metrics(self, results: List[ProcessingResult], batch_time: float):
        """Update performance metrics"""
        self.metrics["total_processed"] += len(results)
        self.metrics["total_processing_time"] += batch_time

        for result in results:
            if result.error:
                self.metrics["errors"] += 1

            file_type = (
                result.file_type.split("/")[0] if "/" in result.file_type else "unknown"
            )
            if file_type not in self.metrics["by_type"]:
                self.metrics["by_type"][file_type] = {"count": 0, "total_time": 0.0}

            self.metrics["by_type"][file_type]["count"] += 1
            self.metrics["by_type"][file_type]["total_time"] += result.processing_time

    def get_performance_metrics(self) -> Dict[str, Any]:
        """Get processing performance metrics"""
        avg_time = self.metrics["total_processing_time"] / max(
            1, self.metrics["total_processed"]
        )

        return {
            "total_processed": self.metrics["total_processed"],
            "total_processing_time": self.metrics["total_processing_time"],
            "average_processing_time": avg_time,
            "error_rate": self.metrics["errors"]
            / max(1, self.metrics["total_processed"]),
            "throughput": self.metrics["total_processed"]
            / max(1, self.metrics["total_processing_time"]),
            "by_type": self.metrics["by_type"],
        }

    def cleanup(self):
        """Clean up resources"""
        self.executor.shutdown(wait=True)
        logger.info("Evidence processor cleaned up")


# Create singleton instance
evidence_processor = EvidenceProcessor(max_workers=4)
